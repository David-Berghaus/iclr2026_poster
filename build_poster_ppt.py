from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/Users/dberghaus/Projects/EVIL_paper/FIM_PointProcess/poster")
ASSETS = ROOT / "ppt_assets"
OUT = ROOT / "fim_pp_poster_recreated.pptx"

WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(36, 41, 49)
MUTED = RGBColor(94, 104, 119)
PANEL = RGBColor(244, 248, 253)
BORDER = RGBColor(208, 220, 235)


def add_box(slide, left, top, width, height, fill=WHITE, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    if radius:
        shape.adjustments[0] = 0.02
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size,
    bold=False,
    color=TEXT,
    italic=False,
    align=PP_ALIGN.LEFT,
    margins=(0.02, 0.01, 0.02, 0.01),
    font_name="Arial",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(margins[0])
    tf.margin_top = Inches(margins[1])
    tf.margin_right = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font_name
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=12):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.01)
    tf.margin_top = Inches(0.01)
    tf.margin_right = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.bullet = True
        p.level = 0
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT
        run.font.name = "Arial"
    return box


def main():
    prs = Presentation()
    prs.slide_width = Inches(18.5)
    prs.slide_height = Inches(9.0)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Start from the LaTeX poster rendering so the PowerPoint matches the source closely.
    slide.shapes.add_picture(
        str(ROOT / "fim_pp_poster.pdf.png"),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

    # Rebalance logos.
    add_box(slide, Inches(11.00), Inches(0.00), Inches(7.10), Inches(1.02))
    logo_specs = [
        (ROOT / "lamarr-logo-2023.png", 11.65, 0.19, 0.46),
        (ROOT / "uni_bonn.png", 13.80, 0.14, 0.55),
        (ROOT / "iais.png", 15.58, 0.14, 0.58),
        (ROOT / "tu-potsdam-flat.png", 17.55, 0.18, 0.42),
    ]
    for path, left, top, h in logo_specs:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), height=Inches(h))

    # Larger QR code while keeping the original title block.
    add_box(slide, Inches(4.40), Inches(0.98), Inches(0.62), Inches(0.48))
    slide.shapes.add_picture(str(ROOT / "qrcode.svg.png"), Inches(4.43), Inches(0.96), height=Inches(0.58))

    # Summary block: larger text and cleaned key message panel.
    add_box(slide, Inches(0.00), Inches(1.84), Inches(5.34), Inches(1.18))
    add_textbox(
        slide,
        Inches(0.10),
        Inches(1.92),
        Inches(5.10),
        Inches(1.00),
        "Modeling irregular marked event sequences as temporal point processes gives an interpretable description of when events occur and which type they have. Existing neural TPP models usually learn one dataset at a time. FIM-PP is a pretrained recognition model that infers conditional intensities in context from a set of related event sequences.\n\nWe pretrain on 72K synthetic processes and 14.4M events. The same model can then be applied zero-shot to real-world event data or adapted within minutes by fine-tuning.",
        font_size=11.0,
    )
    add_box(slide, Inches(0.08), Inches(3.03), Inches(5.08), Inches(0.64), fill=PANEL, line=BORDER, radius=True)
    add_textbox(
        slide,
        Inches(0.16),
        Inches(3.12),
        Inches(4.96),
        Inches(0.46),
        "Key message\nOne pretrained model already competes with specialized TPP baselines; fine-tuning turns that prior into the strongest long-horizon performer.",
        font_size=10.3,
    )

    # Motivation block: consistent spacing and a larger image.
    add_box(slide, Inches(0.00), Inches(4.00), Inches(5.34), Inches(4.58))
    add_textbox(slide, Inches(0.10), Inches(4.10), Inches(2.85), Inches(0.18), "Marked temporal point processes", font_size=10.8, bold=True)
    add_textbox(
        slide,
        Inches(0.10),
        Inches(4.31),
        Inches(2.82),
        Inches(0.86),
        "A marked temporal point process describes an event history Hₜ = {(t₁, k₁), ..., (tₙ, kₙ)} through a conditional intensity λ(t, k | Hₜ) that determines how likely an event of mark k is to occur at time t given the past.",
        font_size=10.1,
    )
    add_textbox(slide, Inches(0.10), Inches(5.22), Inches(2.85), Inches(0.18), "What we want", font_size=10.8, bold=True)
    add_textbox(
        slide,
        Inches(0.10),
        Inches(5.43),
        Inches(2.82),
        Inches(0.66),
        "From a set of observed sequences from one system, infer an interpretable mark-wise intensity model that transfers across datasets without retraining from scratch each time.",
        font_size=10.1,
    )
    add_textbox(slide, Inches(0.10), Inches(6.12), Inches(2.85), Inches(0.18), "Why current practice is limiting", font_size=10.8, bold=True)
    add_bullets(
        slide,
        Inches(0.10),
        Inches(6.31),
        Inches(2.82),
        Inches(0.52),
        [
            "Standard neural TPPs relearn for every dataset.",
            "Zero-shot application is largely unavailable.",
        ],
        font_size=10.0,
    )
    add_textbox(slide, Inches(0.10), Inches(6.88), Inches(2.85), Inches(0.18), "What FIM-PP changes", font_size=10.8, bold=True)
    add_bullets(
        slide,
        Inches(0.10),
        Inches(7.08),
        Inches(2.82),
        Inches(0.80),
        [
            "Pretrain once on a broad synthetic prior.",
            "Infer new dynamics directly from context.",
            "Keep access to an explicit intensity estimate.",
        ],
        font_size=10.0,
    )
    slide.shapes.add_picture(
        str(ROOT / "pp_visualization.jpeg"),
        Inches(3.15),
        Inches(4.40),
        width=Inches(2.02),
        height=Inches(3.10),
    )

    # Part A block: larger equation and body copy.
    add_box(slide, Inches(5.68), Inches(1.86), Inches(6.10), Inches(1.90))
    add_textbox(slide, Inches(5.82), Inches(1.95), Inches(3.80), Inches(0.18), "Part A: Synthetic training data generation", font_size=11.0, bold=True)
    add_textbox(
        slide,
        Inches(5.82),
        Inches(2.15),
        Inches(4.20),
        Inches(0.20),
        "We sample processes from a broad prior over conditional intensities",
        font_size=10.8,
    )
    add_box(slide, Inches(6.48), Inches(2.35), Inches(4.35), Inches(0.62))
    slide.shapes.add_picture(str(ASSETS / "equation.png"), Inches(6.55), Inches(2.43), width=Inches(4.15))
    add_textbox(
        slide,
        Inches(5.88),
        Inches(2.88),
        Inches(5.72),
        Inches(0.70),
        "where μₖ(t) is chosen as either a constant, a positive sinusoid, or a Gamma-shaped initial-rate function, and γₖₖ′(t) is chosen as either zero, an exponential Hawkes kernel, or a shifted Rayleigh kernel. Interaction signs zₖₖ′ ∈ {−1, 0, 1} cover excitatory, inhibitory, and neutral relations.",
        font_size=10.2,
    )

    prs.save(OUT)


if __name__ == "__main__":
    main()
